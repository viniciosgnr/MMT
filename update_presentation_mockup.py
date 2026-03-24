import os
import sys
import subprocess

# Ensure python-pptx is installed in the native environment
try:
    import pptx
except ImportError:
    print("python-pptx missing. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

INPUT_FILE = '/home/marcosgnr/MMT/MMT Specifications/MMT_MVP_MOCKUPS.pptx'

NAVY = RGBColor(0x00, 0x3D, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
SUBTITLE_CLR = RGBColor(0xBB, 0xCC, 0xDD)
MUTED = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

MODULES = {
    "M1": {
        "n": "Equipment Registration & Installation",
        "o": "Manages serial numbers, tags, and equipment lifecycle. Tracks all certificates and installation history for each device.",
        "t": ["FastAPI", "SQLAlchemy", "React Data Tables", "Multi-step Installation Wizard"],
        "f": ["Serial number creation with equipment type and specs", "Tag register with occupancy validation", "Equipment installation and removal history", "Certificate management linked to serial numbers"],
        "d": ["Equipment Change Report (Appendix D template)", "Installation/removal checklist with mandatory steps", "Calibration frequency conflict resolution", "Bulk equipment import from spreadsheets"]
    },
    "M2": {
        "n": "Metrological Confirmation",
        "o": "Manages the full calibration lifecycle per ISO 10012 — from planning through certificate issuance, Critical Analysis, and FC update.",
        "t": ["FastAPI", "IntegrationService (M2→M9)", "Recharts for Uncertainty Visualization"],
        "f": ["Task lifecycle: Pending → Planned → Executed → Completed", "In-situ / ex-situ calibration with seal recording", "Certificate upload and temporary completion dates", "Critical Analysis with 3 hard-limit validation rules", "Uncertainty budget table (ISO 5167)"],
        "d": ["XML auto-population from imported certificates", "Override value calculation per ANP Resolução 18", "UC import and on-demand calculation", "PDF/Excel export in Portuguese and English"]
    },
    "M3": {
        "n": "Chemical Analysis",
        "o": "Manages the sampling workflow from collection through lab validation and flow computer update, with SLA tracking.",
        "t": ["FastAPI", "SLA Matrix Engine", "Critical Analysis dynamic rules"],
        "f": ["Sampling categories (Well Test/Gas, Daily Oil, Periodic)", "11-step SLA pipeline tracking cards", "Automated Critical Analysis (2σ + hard limits)", "Emergency sample scheduling on reproval"],
        "d": ["CRO density extraction and gas composition parsing", "Vendor/warehouse logistics tracking", "Bulk sample creation and import", "Lab report template standardization"]
    },
    "M4": {
        "n": "Onshore Maintenance",
        "o": "Tracks maintenance activities at onshore facilities, linking work orders to serial numbers.",
        "t": ["FastAPI", "SQLAlchemy", "React"],
        "f": ["Maintenance record creation", "Work order tracking linked to serial numbers", "Onshore facility management"],
        "d": ["IFS work order API integration", "Maintenance scheduling and reminders", "Certificate linking to maintenance outcomes", "Spare parts tracking from IFS"]
    },
    "M5": {
        "n": "Synchronize Data",
        "o": "Enables data exchange with external systems (IFS for spare parts, Power BI for reporting).",
        "t": ["FastAPI", "Integration service layer"],
        "f": ["Sync configuration structure", "Sync failure alert generation", "API-first architecture ready"],
        "d": ["IFS API integration for spare parts", "Power BI direct reporting connections", "Scheduled synchronization jobs", "Data consistency mapping dictionaries"]
    },
    "M6": {
        "n": "Monitoring & Alerts",
        "o": "Monitors SLA compliance, FC configuration, and generates proactive alerts for approaching deadlines.",
        "t": ["FastAPI", "SLA Alert Engine", "Recharts"],
        "f": ["Dynamic SLA countdown notifications", "Alert dashboard with severity grouping", "FC Configuration Check vs approved certificates", "Visual budget tables for uncertainty"],
        "d": ["Daily polling of FC parameter changes", "Alarm limit checks across historical data (2-sigma)", "Missing document audits", "Automated email anomaly reports"]
    },
    "M7": {
        "n": "Export Data",
        "o": "Provides structured data export (PDF, Excel, XML) for ANP audits and operational reporting.",
        "t": ["FastAPI", "ZIP archive builder", "Excel/PDF Generators"],
        "f": ["Excel and PDF static format exports", "FPSO selection and date interval filters", "Responsive download modals"],
        "d": ["Strict ANP XML structures (Res. 65/2014 & Ofício 906/2021)", "ZIP packaging matching specific folder tree guidelines", "Bilingual formatting", "Permanent URL generation"]
    },
    "M8": {
        "n": "Planning",
        "o": "Unified view of upcoming calibrations, inspections, and samplings with calendar/list views.",
        "t": ["FastAPI", "React Big Calendar"],
        "f": ["Month/Week/Day timeline views", "List mode with traffic-light status", "Quick task completion popups", "Global filtering by FPSO and dates"],
        "d": ["'Mitigated' status logic for deferred tasks", "Risk forms for date extensions", "Grouping features for overlapping FPSO activities", "Integration with MS Planner/Teams"]
    },
    "M9": {
        "n": "Failure Notification (NFSM)",
        "o": "Automates mismeasurement reports (Res. 18/2014) triggered by calibration failures or system alerts.",
        "t": ["FastAPI", "Approval Workflow Engine"],
        "f": ["Initial (240h) and Final NFSM generation", "Integration directly from M2 calibration failures", "Status flow (Draft → Pending ME → Approved → Sent)", "PDF export"],
        "d": ["ANP XML generation for automatic submissions", "Email distribution lists logic per FPSO", "Tight deadlines tracking (72h rules)", "Calculation interface for production override volumes"]
    },
    "M10": {
        "n": "Historical Reports",
        "o": "Maintains a central, searchable vault of historical PDF documents and audits.",
        "t": ["FastAPI", "File Storage Backend"],
        "f": ["File upload with metadata tagging", "Filtering by FPSO, system, report type", "Audit trail for uploads/downloads"],
        "d": ["Bulk multi-file upload capability", "Optical search/metadata scraping", "Versioning system for replaced reports", "Granular document access permissions"]
    },
    "M11": {
        "n": "Configuration & Asset Hierarchy",
        "o": "Centralizes system parametrization. Models the FPSO hierarchy and manages dynamic properties.",
        "t": ["FastAPI", "SQLAlchemy", "React Tabs", "Recursive Tree Data"],
        "f": ["FPSO tree structure builder", "Dynamic attributes generator (Num, Text, Dates)", "Global admin validation rules mapping", "Wells CRUD with production properties", "Holiday calendars and stock locations"],
        "d": ["Migrating real hierarchy maps from all FPSOs", "Advanced cross-field logic (Rule A depends on Rule B)", "Asset class inheritance definition", "Exporting node configurations"]
    },
    "M12": {
        "n": "User & Access Management",
        "o": "Manages authentication, sessions, and role permissions.",
        "t": ["FastAPI Security", "JWT", "Bcrypt"],
        "f": ["User registration and JWT token login", "Role mapping (Admin, Eng, Tech)", "Password recovery mechanism", "Protected API middleware routing"],
        "d": ["FPSO-level scoped access rules", "SSO integration with Microsoft AD", "User permission administration GUI", "Immutable global access logs view"]
    }
}

def clear_slide(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, lines_data):
    """lines_data: list of (text, font_size, bold, color_rgb)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fsize, bold, color) in enumerate(lines_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(text)
        p.font.size = Pt(fsize)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox

def make_overview(slide, slide_width):
    clear_slide(slide)
    m = Inches(0.5)
    cw = slide_width - 2 * m

    add_rect(slide, 0, 0, slide_width, Inches(1.3), NAVY)
    add_textbox(slide, m, Inches(0.2), cw, Inches(0.5), [
        ("MVP Overview — Metering Management Tool", 24, True, WHITE)
    ])
    add_textbox(slide, m, Inches(0.72), cw, Inches(0.5), [
        ("Web application for managing the complete lifecycle of fiscal and operational metering systems across the SBM Offshore FPSO fleet in Brazil.", 11, False, SUBTITLE_CLR)
    ])

    y = Inches(1.5)
    col_w = Inches(4.5)
    left_lines = [
        ("MVP OBJECTIVE", 10, True, MUTED),
        ("Deliver a functional baseline covering core compliance workflows —", 11, False, DARK),
        ("calibration management, chemical analysis, equipment tracking,", 11, False, DARK),
        ("monitoring/alerts, planning, failure notifications, and historical", 11, False, DARK),
        ("reports — to validate the architecture and gather user feedback.", 11, False, DARK),
        ("", 8, False, DARK),
        ("TECHNOLOGY STACK", 10, True, MUTED),
        ("• Frontend: Next.js 15, React 19, TypeScript, Tailwind CSS, Recharts", 11, False, DARK),
        ("• Backend: Python FastAPI, SQLAlchemy ORM, JWT Auth", 11, False, DARK),
        ("• Database: SQLite (MVP) → PostgreSQL (production)", 11, False, DARK),
        ("• Architecture: REST API-first, SPA + API pattern", 11, False, DARK),
    ]
    add_textbox(slide, m, y, col_w, Inches(3.5), left_lines)

    rx = m + col_w + Inches(0.3)
    rw = cw - col_w - Inches(0.3)
    right_lines = [("MODULES IN THIS MVP", 10, True, MUTED)]
    for mk in ["M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10", "M11", "M12"]:
        right_lines.append((f"{mk} - {MODULES[mk]['n']}", 10, False, DARK))
        
    add_textbox(slide, rx, y, rw, Inches(3.5), right_lines)

def make_module(slide, slide_width, mid, name, obj, tech, feats, todo):
    clear_slide(slide)
    m = Inches(0.5)
    cw = slide_width - 2 * m
    header_h = Inches(1.1)

    add_rect(slide, 0, 0, slide_width, header_h, NAVY)

    add_textbox(slide, m, Inches(0.15), cw, Inches(0.45), [
        (f"{mid} — {name}", 22, True, WHITE)
    ])
    add_textbox(slide, m, Inches(0.6), cw, Inches(0.4), [
        (obj, 11, False, SUBTITLE_CLR)
    ])

    y_pos = header_h + Inches(0.2)
    col_width = Inches(4.5)

    # LEFT
    tech_lines = [("TECHNOLOGIES", 10, True, MUTED)]
    for t in tech:
        tech_lines.append(("• " + t, 11, False, DARK))
    tech_lines.append(("", 5, False, DARK))
    tech_lines.append(("IMPLEMENTED FEATURES", 10, True, GREEN))
    for f in feats:
        tech_lines.append(("✓ " + f, 11, False, DARK))
    add_textbox(slide, m, y_pos, col_width, Inches(4.0), tech_lines)

    # RIGHT
    right_col_x = m + col_width + Inches(0.3)
    right_col_width = cw - col_width - Inches(0.3)
    dev_lines = [("TO BE DEVELOPED / ENHANCED (FINAL VERSION)", 10, True, ORANGE)]
    for d in todo:
        dev_lines.append(("→ " + d, 11, False, DARK))
    add_textbox(slide, right_col_x, y_pos, right_col_width, Inches(4.0), dev_lines)

def main():
    print("Opening presentation...")
    prs = Presentation(INPUT_FILE)
    
    # 1) Find the true blank slides securely
    blank_indices = []
    for i, slide in enumerate(prs.slides):
        texts = []
        imgs = 0
        for s in slide.shapes:
            try:
                if hasattr(s, "image") and getattr(s, "image"): imgs += 1
            except: pass
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    td = p.text.strip()
                    if td: texts.append(td)
        if len(texts) == 0 and imgs == 0:
            blank_indices.append(i)

    print(f"Blank slides found: {blank_indices}")
    
    # 2) Extract ALL text to detect modules mapping easily
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        txt = ""
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    txt += " " + p.text.lower()
        slide_texts.append(txt)

    for i_idx, slide_idx in enumerate(blank_indices):
        slide = prs.slides[slide_idx]
        if i_idx == 0:
            print(f"Slide {slide_idx+1}: Formatting as MVP Overview")
            make_overview(slide, prs.slide_width)
            continue

        # Look ahead 1-3 slides to infer the module
        mod = None
        for j in range(slide_idx + 1, min(slide_idx + 4, len(prs.slides))):
            combined = slide_texts[j]
            mapping = {
                "M11": ["configuration", "asset hierarchy", "fpso tree"],
                "M12": ["user", "access management", "authentication", "login"],
                "M1":  ["equipment", "serial number", "tag register", "multi step installation"],
                "M2":  ["metrological", "calibration"],
                "M3":  ["chemical", "sampling", "analysis", "chromatograph"],
                "M4":  ["maintenance", "onshore"],
                "M5":  ["synchronize", "sync data"],
                "M6":  ["monitoring", "alert", "budget", "uncertainty base"],
                "M7":  ["export"],
                "M8":  ["planning", "calendar"],
                "M9":  ["failure", "nfsm"],
                "M10": ["historical", "report", "documents"]
            }
            for k, keywords in mapping.items():
                if any(kw in combined for kw in keywords):
                    mod = k
                    break
            if mod: break

        if mod and mod in MODULES:
            m_data = MODULES[mod]
            print(f"Slide {slide_idx+1}: Mapping detected -> {mod}")
            make_module(slide, prs.slide_width, mod, m_data["n"], m_data["o"], m_data["t"], m_data["f"], m_data["d"])
        else:
            print(f"Slide {slide_idx+1}: Unidentified module. Left blank.")

    print("Saving presentation in place...")
    prs.save(INPUT_FILE)
    print("Success! PowerPoint updated.")

if __name__ == "__main__":
    main()
